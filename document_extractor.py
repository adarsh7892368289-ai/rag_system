"""
Universal Document Extractor - Production Version

With smart chunking strategies for optimal retrieval performance.
"""

import os
import json
from typing import List, Dict, Optional
from dataclasses import dataclass, asdict
from datetime import datetime
from pathlib import Path
from urllib.parse import urlparse

import scrapy
from scrapy.crawler import CrawlerProcess
import requests
from bs4 import BeautifulSoup
import pdfplumber
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation

# Import chunking module
from chunking import chunk_text


SCRAPY_RESULTS = []


@dataclass
class Document:
    """Standard document structure for embeddings"""
    content: str
    source: str
    doc_id: str
    source_type: str
    metadata: Dict
    chunk_id: Optional[int] = None
    
    def to_dict(self):
        return asdict(self)


class TextSpider(scrapy.Spider):
    """Scrapy spider for web scraping"""
    name = 'text_spider'
    
    def __init__(self, urls=None, follow_links=False, max_pages=10, *args, **kwargs):
        super().__init__(*args, **kwargs)
        
        if isinstance(urls, list):
            self.start_urls = urls
        elif urls:
            self.start_urls = [urls]
        else:
            self.start_urls = []
        
        self.follow_links = follow_links
        self.max_pages = max_pages
        self.pages_scraped = 0
    
    def parse(self, response):
        if self.pages_scraped >= self.max_pages:
            return
        
        self.pages_scraped += 1
        
        title = response.css('title::text').get() or 'No title'
        
        text_parts = []
        for selector in ['p::text', 'h1::text', 'h2::text', 'h3::text', 'li::text', 'td::text']:
            text_parts.extend(response.css(selector).getall())
        
        text = ' '.join([t.strip() for t in text_parts if t.strip()])
        
        SCRAPY_RESULTS.append({
            'content': text,
            'source': response.url,
            'source_type': 'web',
            'metadata': {
                'title': title,
                'extracted_at': datetime.now().isoformat(),
                'char_count': len(text),
                'word_count': len(text.split())
            }
        })
        
        print(f"   ✓ {response.url[:70]}... ({len(text):,} chars)")
        
        if self.follow_links and self.pages_scraped < self.max_pages:
            for link in response.css('a::attr(href)').getall():
                if link and not link.startswith('#') and not link.startswith('javascript:'):
                    yield response.follow(link, callback=self.parse)


class DocumentExtractor:
    """
    Extract and chunk text from various sources
    
    Args:
        chunk_method: 'sentence_aware', 'semantic', 'fixed_size', 'paragraph'
        chunk_size: Target chunk size in words
        chunk_overlap: Overlap between chunks in words
        output_dir: Directory to save extracted documents
    """
    
    def __init__(self, 
                 chunk_method: str = 'sentence_aware',
                 chunk_size: int = 150,
                 chunk_overlap: int = 30,
                 output_dir: str = 'extracted_docs'):
        self.chunk_method = chunk_method
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.output_dir = output_dir
        os.makedirs(self.output_dir, exist_ok=True)
        
        print(f"📋 Chunking method: {chunk_method}")
        print(f"📏 Chunk size: {chunk_size} words (overlap: {chunk_overlap})")
    
    def extract_multiple(self, sources: List[str], chunk: bool = True, 
                        use_scrapy: bool = True, follow_links: bool = False,
                        max_pages: int = 10) -> List[Document]:
        """
        Extract from multiple sources (URLs and files)
        """
        all_documents = []
        
        urls = [s for s in sources if s.startswith(('http://', 'https://'))]
        files = [s for s in sources if not s.startswith(('http://', 'https://'))]
        
        # Process URLs with Scrapy
        if urls and use_scrapy:
            mode = "crawling" if follow_links else "scraping"
            print(f"\n🕷️  Scrapy {mode}: {len(urls)} seed URL(s)")
            if follow_links:
                print(f"   Max pages per seed: {max_pages}")
            
            try:
                raw_docs = self._extract_web_scrapy_batch(urls, follow_links, max_pages)
                all_documents.extend(self._convert_to_documents(raw_docs, chunk))
                        
            except Exception as e:
                print(f"❌ Scrapy failed: {e}")
                print("   Falling back to BeautifulSoup...")
                for url in urls:
                    try:
                        docs = self.extract(url, chunk, use_scrapy=False)
                        all_documents.extend(docs)
                    except Exception as e2:
                        print(f"❌ {url}: {e2}")
        
        # Process files
        if files:
            print(f"\n📁 Processing {len(files)} file(s)...")
        for idx, file_path in enumerate(files, 1):
            print(f"   [{idx}/{len(files)}] {Path(file_path).name}")
            try:
                docs = self.extract(file_path, chunk, use_scrapy=False)
                all_documents.extend(docs)
            except Exception as e:
                print(f"   ❌ Error: {e}")
        
        return all_documents
    
    def extract(self, source: str, chunk: bool = True, use_scrapy: bool = False,
                follow_links: bool = False, max_pages: int = 10) -> List[Document]:
        """Extract from a single source"""
        if source.startswith(('http://', 'https://')):
            if use_scrapy:
                raw_docs = self._extract_web_scrapy_batch([source], follow_links, max_pages)
            else:
                raw_docs = self._extract_web_bs4(source)
        else:
            if not os.path.exists(source):
                raise FileNotFoundError(f"File not found: {source}")
            
            ext = Path(source).suffix.lower()
            extractors = {
                '.pdf': self._extract_pdf,
                '.xlsx': self._extract_excel,
                '.xls': self._extract_excel,
                '.csv': self._extract_csv,
                '.docx': self._extract_word,
                '.pptx': self._extract_pptx,
                '.txt': self._extract_txt,
                '.json': self._extract_json,
            }
            
            if ext not in extractors:
                raise ValueError(f"Unsupported file type: {ext}")
            
            raw_docs = extractors[ext](source)
        
        return self._convert_to_documents(raw_docs, chunk)
    
    def _convert_to_documents(self, raw_docs: List[Dict], do_chunk: bool) -> List[Document]:
        """Convert raw extracted data to Document objects with optional chunking"""
        documents = []
        
        for raw_doc in raw_docs:
            if do_chunk:
                chunks = self._chunk_text(raw_doc['content'])
                for idx, chunk_text in enumerate(chunks):
                    chunk_metadata = {
                        'chunk_char_count': len(chunk_text),
                        'chunk_word_count': len(chunk_text.split()),
                        **raw_doc['metadata']
                    }
                    
                    doc = Document(
                        content=chunk_text,
                        source=raw_doc['source'],
                        doc_id=f"{raw_doc['source']}_{idx}",
                        source_type=raw_doc['source_type'],
                        metadata=chunk_metadata,
                        chunk_id=idx
                    )
                    documents.append(doc)
            else:
                doc = Document(
                    content=raw_doc['content'],
                    source=raw_doc['source'],
                    doc_id=raw_doc['source'],
                    source_type=raw_doc['source_type'],
                    metadata=raw_doc['metadata'],
                    chunk_id=None
                )
                documents.append(doc)
        
        return documents
    
    def _chunk_text(self, text: str) -> List[str]:
        """Chunk text using configured method"""
        return chunk_text(
            text,
            method=self.chunk_method,
            target_words=self.chunk_size,
            overlap_words=self.chunk_overlap
        )
    
    def _extract_web_scrapy_batch(self, urls: List[str], follow_links: bool = False,
                                   max_pages: int = 10) -> List[Dict]:
        """Extract from URLs using Scrapy"""
        global SCRAPY_RESULTS
        SCRAPY_RESULTS = []
        
        process = CrawlerProcess(settings={
            'LOG_LEVEL': 'ERROR',
            'ROBOTSTXT_OBEY': True,
            'CONCURRENT_REQUESTS': 2,
            'DOWNLOAD_DELAY': 1,
            'USER_AGENT': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36',
            'DOWNLOAD_TIMEOUT': 30,
            'RETRY_TIMES': 2,
        })
        
        process.crawl(TextSpider, urls=urls, follow_links=follow_links, max_pages=max_pages)
        process.start()
        
        return SCRAPY_RESULTS.copy()
    
    def _extract_web_bs4(self, url: str) -> List[Dict]:
        """Extract from URL using BeautifulSoup"""
        try:
            headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
            response = requests.get(url, headers=headers, timeout=30)
            response.raise_for_status()
            
            soup = BeautifulSoup(response.content, 'lxml')
            title = soup.find('title')
            title = title.text.strip() if title else 'No title'
            
            for element in soup(['script', 'style', 'nav', 'footer', 'header', 'aside']):
                element.decompose()
            
            text = soup.get_text(separator=' ', strip=True)
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = ' '.join(chunk for chunk in chunks if chunk)
            
            print(f"   ✓ {url[:70]}... ({len(text):,} chars)")
            
            return [{
                'content': text,
                'source': url,
                'source_type': 'web',
                'metadata': {
                    'title': title,
                    'extracted_at': datetime.now().isoformat(),
                    'char_count': len(text),
                    'word_count': len(text.split())
                }
            }]
            
        except Exception as e:
            print(f"   ❌ Error: {str(e)[:50]}")
            return []
    
    def _extract_pdf(self, file_path: str) -> List[Dict]:
        """Extract text from PDF"""
        with pdfplumber.open(file_path) as pdf:
            full_text = []
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    full_text.append(text)
            
            combined_text = '\n\n'.join(full_text)
            
            return [{
                'content': combined_text,
                'source': file_path,
                'source_type': 'pdf',
                'metadata': {
                    'total_pages': len(pdf.pages),
                    'extracted_at': datetime.now().isoformat(),
                    'char_count': len(combined_text),
                    'word_count': len(combined_text.split())
                }
            }]
    
    def _extract_excel(self, file_path: str) -> List[Dict]:
        """Extract text from Excel"""
        documents = []
        excel_file = pd.ExcelFile(file_path)
        
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            
            text_lines = [' | '.join(str(col) for col in df.columns)]
            for _, row in df.iterrows():
                text_lines.append(' | '.join(str(val) for val in row.values))
            
            text = '\n'.join(text_lines)
            
            documents.append({
                'content': text,
                'source': f"{file_path}::{sheet_name}",
                'source_type': 'excel',
                'metadata': {
                    'sheet_name': sheet_name,
                    'rows': len(df),
                    'columns': len(df.columns),
                    'extracted_at': datetime.now().isoformat()
                }
            })
        
        return documents
    
    def _extract_csv(self, file_path: str) -> List[Dict]:
        """Extract text from CSV"""
        df = pd.read_csv(file_path)
        
        text_lines = [' | '.join(str(col) for col in df.columns)]
        for _, row in df.iterrows():
            text_lines.append(' | '.join(str(val) for val in row.values))
        
        text = '\n'.join(text_lines)
        
        return [{
            'content': text,
            'source': file_path,
            'source_type': 'csv',
            'metadata': {
                'rows': len(df),
                'columns': len(df.columns),
                'extracted_at': datetime.now().isoformat()
            }
        }]
    
    def _extract_word(self, file_path: str) -> List[Dict]:
        """Extract text from Word document"""
        doc = DocxDocument(file_path)
        paragraphs = [para.text for para in doc.paragraphs if para.text.strip()]
        
        tables_text = []
        for table in doc.tables:
            for row in table.rows:
                tables_text.append(' | '.join(cell.text for cell in row.cells))
        
        text = '\n\n'.join(paragraphs)
        if tables_text:
            text += '\n\n' + '\n'.join(tables_text)
        
        return [{
            'content': text,
            'source': file_path,
            'source_type': 'word',
            'metadata': {
                'paragraphs': len(paragraphs),
                'tables': len(doc.tables),
                'extracted_at': datetime.now().isoformat()
            }
        }]
    
    def _extract_pptx(self, file_path: str) -> List[Dict]:
        """Extract text from PowerPoint"""
        prs = Presentation(file_path)
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_content = getattr(shape, "text", "")
                    if text_content and text_content.strip():
                        slide_text.append(text_content)
            
            if slide_text:
                all_text.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
        
        text = '\n\n'.join(all_text)
        
        return [{
            'content': text,
            'source': file_path,
            'source_type': 'powerpoint',
            'metadata': {
                'total_slides': len(prs.slides),
                'extracted_at': datetime.now().isoformat()
            }
        }]
    
    def _extract_txt(self, file_path: str) -> List[Dict]:
        """Extract text from plain text file"""
        with open(file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        
        return [{
            'content': text,
            'source': file_path,
            'source_type': 'text',
            'metadata': {'extracted_at': datetime.now().isoformat()}
        }]
    
    def _extract_json(self, file_path: str) -> List[Dict]:
        """Extract text from JSON file"""
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        
        text = json.dumps(data, indent=2)
        
        return [{
            'content': text,
            'source': file_path,
            'source_type': 'json',
            'metadata': {'extracted_at': datetime.now().isoformat()}
        }]
    
    def save_documents(self, documents: List[Document]):
        """Save documents grouped by source as JSON files"""
        if not documents:
            print("\n⚠️  No documents to save")
            return
        
        print(f"\n💾 Saving to '{self.output_dir}/'")
        
        docs_by_source = {}
        for doc in documents:
            if doc.source not in docs_by_source:
                docs_by_source[doc.source] = []
            docs_by_source[doc.source].append(doc)
        
        saved_files = []
        for source, docs in docs_by_source.items():
            filename = self._generate_filename(source, docs[0].source_type)
            filepath = os.path.join(self.output_dir, filename)
            
            file_data = {
                'source': source,
                'source_type': docs[0].source_type,
                'chunk_method': self.chunk_method,
                'total_chunks': len(docs),
                'extracted_at': datetime.now().isoformat(),
                'overall_metadata': {
                    'total_char_count': sum(len(doc.content) for doc in docs),
                    'total_word_count': sum(len(doc.content.split()) for doc in docs)
                },
                'chunks': [doc.to_dict() for doc in docs]
            }
            
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(file_data, f, indent=2, ensure_ascii=False)
            
            saved_files.append(filename)
            print(f"   ✓ {filename} ({len(docs)} chunks)")
        
        print(f"\n✅ Saved {len(saved_files)} file(s)")
        self._create_index_file(docs_by_source)
    
    def _generate_filename(self, source: str, source_type: str) -> str:
        """Generate clean filename from source"""
        if source.startswith(('http://', 'https://')):
            parsed = urlparse(source)
            domain = parsed.netloc.replace('www.', '')
            path = parsed.path.strip('/').replace('/', '_')
            base_name = f"{domain}_{path}" if path else domain
        else:
            base_name = Path(source).stem
            if '::' in base_name:
                base_name = base_name.replace('::', '_')
        
        base_name = base_name.replace(' ', '_').replace(':', '_')
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        return f"{source_type}_{base_name}_{timestamp}.json"
    
    def _create_index_file(self, docs_by_source: Dict):
        """Create index file with extraction summary"""
        index_data = {
            'chunk_method': self.chunk_method,
            'chunk_size': self.chunk_size,
            'chunk_overlap': self.chunk_overlap,
            'total_sources': len(docs_by_source),
            'total_chunks': sum(len(docs) for docs in docs_by_source.values()),
            'created_at': datetime.now().isoformat(),
            'sources': []
        }
        
        for source, docs in docs_by_source.items():
            index_data['sources'].append({
                'source': source,
                'source_type': docs[0].source_type,
                'total_chunks': len(docs),
                'total_chars': sum(len(doc.content) for doc in docs),
                'total_words': sum(len(doc.content.split()) for doc in docs),
                'metadata': docs[0].metadata
            })
        
        index_path = os.path.join(self.output_dir, '_index.json')
        with open(index_path, 'w', encoding='utf-8') as f:
            json.dump(index_data, f, indent=2, ensure_ascii=False)


if __name__ == '__main__':
    extractor = DocumentExtractor(
        chunk_method='sentence_aware',
        chunk_size=150,
        chunk_overlap=30,
        output_dir='extracted_docs'
    )
    
    sources = [
        "https://en.wikipedia.org/wiki/Embedding_(machine_learning)",
        "https://www.ibm.com/think/topics/embedding",
        "https://en.wikipedia.org/wiki/BERT_(language_model)",
        "https://en.wikipedia.org/wiki/Word2vec",
        "https://sbert.net/",
        "https://huggingface.co/sentence-transformers",
        "https://en.wikipedia.org/wiki/Machine_learning",
        "https://www.ibm.com/think/topics/machine-learning",
        "https://aws.amazon.com/what-is/embeddings-in-machine-learning/",
        "https://developers.google.com/machine-learning/crash-course/embeddings/embedding-space",
        "https://www.cloudflare.com/learning/ai/what-are-embeddings/"
    ]
    
    print("="*70)
    print("Document Extractor - Production Version")
    print("="*70)
    
    all_docs = extractor.extract_multiple(
        sources, 
        chunk=True, 
        use_scrapy=True,
        follow_links=False
    )
    
    extractor.save_documents(all_docs)
    
    print("\n" + "="*70)
    print(f"✅ Extraction complete: {len(all_docs)} chunks from {len(sources)} sources")
    print(f"📁 Output: {extractor.output_dir}/")
    print("="*70)