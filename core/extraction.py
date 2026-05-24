"""Document extraction pipeline.

Reads heterogeneous sources (web pages, PDFs, docs, spreadsheets, etc.) and
turns each into a set of strategy-tagged chunks ready for indexing. Web and
file extraction are dispatched in a thread pool because they're I/O bound.
"""

import hashlib
import json
import logging
import os
from concurrent.futures import ThreadPoolExecutor, as_completed
from datetime import datetime
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional
from urllib.parse import urlparse

import requests
from bs4 import BeautifulSoup

from config.settings import EXTRACTION
from core.chunking import EnsembleChunker


logger = logging.getLogger(__name__)


_CHUNKING_STRATEGIES = ('sentence_aware', 'semantic', 'paragraph', 'fixed_size')


def _empty_strategy_dict() -> Dict[str, List[Dict[str, Any]]]:
    return {strategy: [] for strategy in _CHUNKING_STRATEGIES}


class DocumentExtractor:

    def __init__(self):
        self._config = EXTRACTION
        self._save_dir = self._config.output_dir
        self._chunker = EnsembleChunker()
        os.makedirs(self._save_dir, exist_ok=True)

    def extract_multiple(
        self, sources: List[str]
    ) -> Dict[str, List[Dict[str, Any]]]:
        """Extract many sources concurrently and merge into per-strategy lists."""
        all_strategy_documents = _empty_strategy_dict()

        print(f"\n📥 Processing {len(sources)} source(s) in parallel...")

        with ThreadPoolExecutor(max_workers=self._config.max_workers) as executor:
            future_to_source = {
                executor.submit(self._extract_one, source): source
                for source in sources
            }

            for future in as_completed(future_to_source):
                source = future_to_source[future]
                try:
                    strategy_docs = future.result()
                    for strategy, docs in strategy_docs.items():
                        all_strategy_documents[strategy].extend(docs)
                    total = sum(len(docs) for docs in strategy_docs.values())
                    print(f"   ✓ {self._display_name(source)}: {total} chunks")
                except Exception as e:
                    print(f"   ✗ {self._display_name(source)}: {e}")

        total_docs = sum(len(docs) for docs in all_strategy_documents.values())
        print(f"\n✅ Extraction complete: {total_docs} total documents across "
              f"{len(_CHUNKING_STRATEGIES)} strategies")
        for strategy, docs in all_strategy_documents.items():
            print(f"   • {strategy}: {len(docs)} chunks")

        return all_strategy_documents

    def save_documents(
        self, strategy_documents: Dict[str, List[Dict[str, Any]]]
    ):
        """Persist extracted chunks to disk, partitioned by source and strategy."""
        if not any(strategy_documents.values()):
            return

        print(f"\n💾 Saving to '{self._save_dir}'")

        for strategy, documents in strategy_documents.items():
            if not documents:
                continue

            by_source: Dict[str, List[Dict]] = {}
            for doc in documents:
                source = doc['metadata']['source']
                by_source.setdefault(source, []).append(doc)

            for source, docs in by_source.items():
                self._save_source_chunks(source, strategy, docs)

    def load_from_folder(
        self, folder_path: str
    ) -> Dict[str, List[Dict[str, Any]]]:
        """Reload previously-saved chunks. Returns the canonical strategy-keyed shape."""
        all_strategy_documents = _empty_strategy_dict()
        folder = Path(folder_path)

        if not folder.exists():
            print(f"⚠️  Folder not found: {folder_path}")
            return all_strategy_documents

        json_files = list(folder.glob('*.json'))
        if not json_files:
            print(f"⚠️  No JSON files in: {folder_path}")
            return all_strategy_documents

        for json_file in json_files:
            try:
                with open(json_file, 'r', encoding='utf-8') as f:
                    data = json.load(f)

                strategy = data.get('chunking_strategy', 'sentence_aware')
                if strategy not in all_strategy_documents:
                    logger.warning(
                        "Skipping %s: unknown strategy %s", json_file.name, strategy
                    )
                    continue

                source = data.get('source', '')
                source_id = self._generate_id(source)
                for i, chunk_data in enumerate(data.get('chunks', [])):
                    all_strategy_documents[strategy].append({
                        'id': f"{source_id}_{strategy}_{i}",
                        'content': chunk_data.get('content', ''),
                        'metadata': chunk_data.get('metadata', {}),
                    })
            except Exception as e:
                print(f"⚠️  Error loading {json_file.name}: {e}")

        return all_strategy_documents

    # ------------------------------------------------------------------ private

    def _extract_one(self, source: str) -> Dict[str, List[Dict[str, Any]]]:
        if source.startswith(('http://', 'https://')):
            return self._extract_web(source)
        return self._extract_file(source)

    def _extract_web(self, url: str) -> Dict[str, List[Dict[str, Any]]]:
        headers = {'User-Agent': self._config.user_agent}
        response = requests.get(url, headers=headers, timeout=self._config.timeout)
        response.raise_for_status()

        soup = BeautifulSoup(response.text, 'html.parser')
        # Strip non-content elements before extracting text — otherwise nav/footer
        # boilerplate dominates the chunks.
        for tag in soup(['script', 'style', 'nav', 'header', 'footer', 'aside']):
            tag.decompose()

        title_tag = soup.find('title')
        title_text = title_tag.get_text().strip() if title_tag else url

        content = ' '.join(soup.get_text(separator=' ', strip=True).split())
        if len(content) < self._config.min_web_content_chars:
            raise ValueError(
                f"Insufficient content extracted ({len(content)} chars)"
            )

        base_metadata = {
            'source': url,
            'source_type': 'web',
            'title': title_text,
            'url': url,
            'domain': urlparse(url).netloc,
            'extracted_at': datetime.now().isoformat(),
        }
        return self._chunk_with_all_strategies(content, base_metadata)

    def _extract_file(self, file_path: str) -> Dict[str, List[Dict[str, Any]]]:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = path.suffix.lower()
        extractor = self._file_extractors().get(ext)
        if extractor is None:
            raise ValueError(f"Unsupported file type: {ext}")

        content = extractor(path)

        base_metadata = {
            'source': str(path),
            'source_type': ext.lstrip('.'),
            'title': path.stem,
            'filename': path.name,
            'file_type': ext,
            'file_size': path.stat().st_size,
            'extracted_at': datetime.now().isoformat(),
        }
        return self._chunk_with_all_strategies(content, base_metadata)

    def _file_extractors(self) -> Dict[str, Callable[[Path], str]]:
        return {
            '.pdf': self._extract_pdf,
            '.docx': self._extract_word,
            '.doc': self._extract_word,
            '.xlsx': self._extract_excel,
            '.xls': self._extract_excel,
            '.csv': self._extract_csv,
            '.pptx': self._extract_pptx,
            '.txt': self._extract_txt,
            '.json': self._extract_json,
        }

    def _extract_pdf(self, path: Path) -> str:
        try:
            import pdfplumber
        except ImportError as e:
            raise ImportError("pdfplumber required: pip install pdfplumber") from e

        with pdfplumber.open(path) as pdf:
            pages = [page.extract_text() for page in pdf.pages]
        return '\n\n'.join(p for p in pages if p)

    def _extract_word(self, path: Path) -> str:
        try:
            from docx import Document
        except ImportError as e:
            raise ImportError("python-docx required: pip install python-docx") from e

        doc = Document(path)
        return '\n\n'.join(p.text for p in doc.paragraphs if p.text.strip())

    def _extract_excel(self, path: Path) -> str:
        try:
            import pandas as pd
        except ImportError as e:
            raise ImportError(
                "pandas and openpyxl required: pip install pandas openpyxl"
            ) from e

        excel_file = pd.ExcelFile(path)
        sheets = []
        for sheet_name in excel_file.sheet_names:
            df = pd.read_excel(excel_file, sheet_name=sheet_name)
            lines = [
                f"Sheet: {sheet_name}",
                ' | '.join(str(col) for col in df.columns),
            ]
            lines.extend(
                ' | '.join(str(val) for val in row.values)
                for _, row in df.iterrows()
            )
            sheets.append('\n'.join(lines))
        return '\n\n'.join(sheets)

    def _extract_csv(self, path: Path) -> str:
        try:
            import pandas as pd
        except ImportError as e:
            raise ImportError("pandas required: pip install pandas") from e

        df = pd.read_csv(path)
        lines = [' | '.join(str(col) for col in df.columns)]
        lines.extend(
            ' | '.join(str(val) for val in row.values) for _, row in df.iterrows()
        )
        return '\n'.join(lines)

    def _extract_pptx(self, path: Path) -> str:
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ImportError("python-pptx required: pip install python-pptx") from e

        prs = Presentation(path)
        slides = []
        for slide_num, slide in enumerate(prs.slides, start=1):
            shape_texts = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, 'text') and shape.text.strip()
            ]
            if shape_texts:
                slides.append(f"Slide {slide_num}:\n" + '\n'.join(shape_texts))
        return '\n\n'.join(slides)

    def _extract_txt(self, path: Path) -> str:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()

    def _extract_json(self, path: Path) -> str:
        with open(path, 'r', encoding='utf-8') as f:
            data = json.load(f)
        return json.dumps(data, indent=2)

    def _chunk_with_all_strategies(
        self, text: str, base_metadata: Dict[str, Any]
    ) -> Dict[str, List[Dict[str, Any]]]:
        strategy_chunks = self._chunker.chunk_all_strategies(text)
        source_id = self._generate_id(base_metadata['source'])

        strategy_documents: Dict[str, List[Dict[str, Any]]] = {}
        for strategy, chunks in strategy_chunks.items():
            documents = []
            for chunk in chunks:
                documents.append({
                    'id': f"{source_id}_{strategy}_{chunk.index}",
                    'content': chunk.text,
                    'metadata': {
                        **base_metadata,
                        'chunk_index': chunk.index,
                        'chunk_word_count': chunk.word_count,
                        'chunk_char_count': chunk.char_count,
                        **chunk.metadata,
                    },
                })
            strategy_documents[strategy] = documents
        return strategy_documents

    def _save_source_chunks(self, source: str, strategy: str, docs: List[Dict]):
        source_type = docs[0]['metadata']['source_type']
        source_id = self._generate_id(source)
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        base_name = self._safe_basename(source)
        filename = f"{source_type}_{base_name}_{strategy}_{timestamp}_{source_id}.json"
        filepath = os.path.join(self._save_dir, filename)

        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump({
                'source': source,
                'source_type': source_type,
                'chunking_strategy': strategy,
                'total_chunks': len(docs),
                'extracted_at': docs[0]['metadata'].get('extracted_at'),
                'chunks': [
                    {'content': doc['content'], 'metadata': doc['metadata']}
                    for doc in docs
                ],
            }, f, indent=2, ensure_ascii=False)

        print(f"   ✓ {filename} ({len(docs)} chunks)")

    @staticmethod
    def _safe_basename(source: str) -> str:
        if source.startswith(('http://', 'https://')):
            parsed = urlparse(source)
            domain = parsed.netloc.replace('www.', '')
            path_parts = [p for p in parsed.path.split('/') if p]
            page_id = path_parts[-1] if path_parts else 'index'
            page_id = page_id.replace('.html', '').replace('.php', '')[:30]
            base_name = f"{domain}_{page_id}" if page_id != 'index' else domain
        else:
            base_name = Path(source).stem

        return ''.join(c if c.isalnum() or c in '._-' else '_' for c in base_name)

    @staticmethod
    def _display_name(source: str) -> str:
        if source.startswith(('http://', 'https://')):
            return urlparse(source).netloc or source
        return Path(source).name

    @staticmethod
    def _generate_id(source: str) -> str:
        return hashlib.md5(source.encode('utf-8')).hexdigest()[:12]
